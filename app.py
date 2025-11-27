import pandas as pd
import numpy as np
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score, confusion_matrix, mean_squared_error
from sklearn.ensemble import HistGradientBoostingRegressor
from xgboost import XGBClassifier
from sklearn.cluster import KMeans
import streamlit as st
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go
from prophet import Prophet
from openai import OpenAI
import groq
import datetime
from backend.chatbot import get_chat_response
import PyPDF2
import docx
import base64



st.set_page_config(page_title="Crime Dashboard", page_icon="🤖")

st.markdown("""
<style>
[data-testid="stAppViewContainer"] {
    background: linear-gradient(270deg, #000000, #4B0000, #660000, #000000);
    background-size: 800% 800%;
    animation: moveGradient 15s ease infinite;
}

@keyframes moveGradient {
    0% { background-position: 0% 50%; }
    50% { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
}
</style>
""", unsafe_allow_html=True)

st.markdown(
    """
    <style>
    [data-testid="stAppViewContainer"] {
        background-color: #000000;
        color: #ffffff;  /* optional: make text white */
    }
    </style>
    """,
    unsafe_allow_html=True
)

def add_crime_tape():
    tape_path = "images/crime_tape.jpeg"

    # Load and encode tape image
    with open(tape_path, "rb") as f:
        tape_base64 = base64.b64encode(f.read()).decode()

    st.markdown(f"""
        <style>

        /* ---- Animation (slight wiggle effect) ---- */
        @keyframes tapeMove {{
            0%   {{ transform: rotate(-35deg) translateX(0px); }}
            50%  {{ transform: rotate(-35deg) translateX(10px); }}
            100% {{ transform: rotate(-35deg) translateX(0px); }}
        }}

        .crime-tape-top {{
            position: fixed;
            top: 0;
            left: 285px;
            width: 200px;
            z-index: 9999;
            opacity: 0.9;
            transform: rotate(-25deg) skewX(-10deg);
            animation: tapeMove 3s ease-in-out infinite;
            pointer-events: none;
        }}

        .crime-tape-bottom {{
            position: fixed;
            bottom: 0;
            right: -30px;
            width: 175px;
            z-index: 9999;
            opacity: 0.9;
            transform: rotate(25deg) skewX(10deg);
            animation: tapeMove 3s ease-in-out infinite;
            pointer-events: none;
        }}

        </style>

        <img src="data:image/jpeg;base64,{tape_base64}" class="crime-tape-top">
        <img src="data:image/jpeg;base64,{tape_base64}" class="crime-tape-bottom">
    """, unsafe_allow_html=True)

add_crime_tape()


if "chat_started" not in st.session_state:
    st.session_state.chat_started = False

if "messages" not in st.session_state:
    st.session_state.messages = []
    
if "show_chat" not in st.session_state:
    st.session_state.show_chat = False

start_chat = st.sidebar.button("Start Crime AI Assistant")

st.markdown("""
<style>
div.stButton > button:first-child {
    background-color: #ff4b4b;
    color: white;
    border-radius: 10px;
    height: 45px;
    border: 2px solid white;
    font-size: 17px;
    font-weight: 600;
    box-shadow: 0px 0px 12px rgba(255, 75, 75, 0.8);
    transition: 0.3s;
}
div.stButton > button:first-child:hover {
    background-color: #ff2525;
    box-shadow: 0px 0px 20px rgba(255, 0, 0, 1);
}
</style>
""", unsafe_allow_html=True)

if start_chat:
    st.session_state.chat_started = True

    st.session_state.messages.append({
        "role": "assistant",
        "content": "👋 Hello! I am Crime AI Assistant. I analyze crime trends, hotspots, and patterns.\nAsk me anything!"
    })


if "open_modal" not in st.session_state:
    st.session_state.open_modal = False

if "file_text" not in st.session_state:
    st.session_state.file_text = ""

# Sidebar Button to open popup
with st.sidebar:
    if st.button("📄 Upload Crime Report"):
        st.session_state.open_modal = True


# --------------------------------------------------------
# HTML POPUP + CSS (empty box)
# --------------------------------------------------------
if st.session_state.open_modal:

    st.markdown("""
        <style>
        #crime-modal {
            position: fixed;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            width: 480px;
            background: #111;
            color: white;
            padding: 25px;
            border-radius: 18px;
            z-index: 999999;
            box-shadow: 0px 4px 20px rgba(0,0,0,0.6);
        }

        #close-modal {
            position: absolute;
            top: 12px;
            right: 18px;
            font-size: 22px;
            color: #ff4d4d;
            cursor: pointer;
        }

        /* This is the magic part — attach Streamlit widgets INSIDE popup */
        .modal-container {
            position: fixed;
            top: 58%;
            left: 50%;
            transform: translate(-50%, -50%);
            width: 430px;
            z-index: 1000000;
            padding: 10px;
            background: transparent;
        }
        </style>

        <script>
            document.getElementById("close-modal")
                .onclick = function() {
                    window.parent.postMessage({type: "close_popup"}, "*");
                };
        </script>
    """, unsafe_allow_html=True)

    # LISTEN TO JS CLOSE EVENT
    msg = st.experimental_get_query_params()
    if "close_popup" in msg:
        st.session_state.open_modal = False


    # --------------------------------------------------------
    # NOW PLACE STREAMLIT WIDGETS INSIDE POPUP USING CONTAINER
    # --------------------------------------------------------
    with st.container():
        st.markdown('<div class="modal-container">', unsafe_allow_html=True)

        # File Upload
        uploaded_file = st.file_uploader("Choose a PDF or DOCX", type=["pdf", "docx"])

        if uploaded_file:
            if uploaded_file.type == "application/pdf":
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                st.session_state.file_text = "".join(
                    [page.extract_text() or "" for page in pdf_reader.pages]
                )
            else:
                doc_file = docx.Document(uploaded_file)
                st.session_state.file_text = "\n".join([p.text for p in doc_file.paragraphs])

            st.success("File uploaded!")

            # Question Box
            question = st.text_input("Ask a question:")

            # Ask AI Button
            if st.button("Ask AI"):
                if question.strip():
                    final_prompt = (
                        f"Document Content:\n{st.session_state.file_text}\n\n"
                        f"Question: {question}"
                    )
                    ai_answer = get_chat_response(final_prompt)
                    st.write("### 🔍 AI Answer:")
                    st.write(ai_answer)

        st.markdown('</div>', unsafe_allow_html=True)
        
        
        
# -----------------------
# Load API Keys from secrets.toml
# -----------------------
OPENAI_API_KEY = st.secrets["OPENAI_API_KEY"]
GROQ_API_KEY = st.secrets["GROQ_API_KEY"]



# -----------------------
# 1. Load Dataset
# -----------------------
st.title("Crime Pattern Detection and Predictive Analytics Dashboard")
data = pd.read_csv(r"C:\Users\csc\BDA_Project\crime_dataset_india.csv")
st.subheader("Raw Data")
st.dataframe(data.head())

# -----------------------
# 2. Preprocessing
# -----------------------
# Fill missing values
for col in data.select_dtypes(include='object').columns:
    data[col].fillna(data[col].mode()[0], inplace=True)
for col in data.select_dtypes(include='number').columns:
    data[col].fillna(data[col].median(), inplace=True)

# Convert dates safely
for col in ['Date of Occurrence', 'Time of Occurrence', 'Date Case Closed']:
    if col in data.columns:
        data[col] = pd.to_datetime(data[col], errors='coerce', dayfirst=True)

# Feature engineering
data['Hour'] = pd.to_datetime(data['Time of Occurrence'], errors='coerce').dt.hour
data['Month'] = data['Date of Occurrence'].dt.month

# Ensure Victim Age is numeric
if 'Victim Age' in data.columns:
    data['Victim Age'] = pd.to_numeric(data['Victim Age'], errors='coerce').fillna(data['Victim Age'].median())

# Label Encoding for categorical columns (only if present)
cat_cols = ['City', 'Crime Code', 'Crime Description', 'Victim Gender', 'Weapon Used', 'Crime Domain']
le_dict = {}
for col in cat_cols:
    if col in data.columns:
        le = LabelEncoder()
        data[col + '_encoded'] = le.fit_transform(data[col].astype(str))
        le_dict[col] = le

# Preserve original city names for visualization
if 'City_encoded' in data.columns and 'City' in data.columns:
    city_mapping = dict(zip(data['City_encoded'], data['City']))

# -----------------------
# 3. Sidebar Filters
# -----------------------
st.sidebar.title("Filter Options")

month_options = sorted(data['Month'].dropna().unique())
default_months = month_options.copy()

# -----------------------
# City filter with "Select All"
# -----------------------
all_cities = sorted(data['City'].unique()) if 'City' in data.columns else []
select_all_cities = st.sidebar.checkbox("Select All Cities", value=True)

if select_all_cities:
    cities_selected = all_cities
else:
    cities_selected = st.sidebar.multiselect(
        "Select Cities", options=all_cities, default=[all_cities[0]] if all_cities else []
    )

city_codes_selected = []
if 'City' in data.columns and 'City' in le_dict:
    city_codes_selected = [le_dict['City'].transform([c])[0] for c in cities_selected]
elif 'City' in data.columns:
    # fallback if encoding absent
    city_codes_selected = cities_selected

# -----------------------
# Crime domain filter with "Select All"
# -----------------------
all_domains = sorted(data['Crime Domain'].unique()) if 'Crime Domain' in data.columns else []
select_all_domains = st.sidebar.checkbox("Select All Crime Domains", value=True)

if select_all_domains:
    crime_domains_selected = all_domains
else:
    crime_domains_selected = st.sidebar.multiselect(
        "Select Crime Domain", options=all_domains, default=[all_domains[0]] if all_domains else []
    )

crime_codes_selected = []
if 'Crime Domain' in data.columns and 'Crime Domain' in le_dict:
    crime_codes_selected = [le_dict['Crime Domain'].transform([c])[0] for c in crime_domains_selected]
elif 'Crime Domain' in data.columns:
    crime_codes_selected = crime_domains_selected

# -----------------------
# Month filter with "Select All"
# -----------------------
all_months = sorted(data['Month'].dropna().unique()) if 'Month' in data.columns else []
select_all_months = st.sidebar.checkbox("Select All Months", value=True)

if select_all_months:
    months_selected = all_months
else:
    months_selected = st.sidebar.multiselect(
        "Select Month", options=all_months, default=[all_months[0]] if all_months else []
    )

# -----------------------
# Apply filters
# -----------------------
filtered_data = data.copy()
filters = []
if 'City_encoded' in data.columns:
    filtered_data = filtered_data[filtered_data['City_encoded'].isin(city_codes_selected)]
if 'Crime Domain_encoded' in data.columns:
    filtered_data = filtered_data[filtered_data['Crime Domain_encoded'].isin(crime_codes_selected)]
if 'Month' in data.columns and months_selected:
    filtered_data = filtered_data[filtered_data['Month'].isin(months_selected)]

st.subheader("Filtered Data")
st.dataframe(filtered_data.head())

# -----------------------
# NEW: Crime Severity Index Calculation
# -----------------------
# We'll compute a severity score (roughly 1-10) per row using weighted factors.
def calculate_severity(row, city_counts, mean_city_count):
    score = 0.0

    # Weapon factor (High weight)
    weapon = str(row.get('Weapon Used', '')).upper()
    if 'GUN' in weapon or 'FIREARM' in weapon:
        score += 3.0
    elif 'KNIFE' in weapon or 'BLADE' in weapon:
        score += 2.0
    elif weapon in ('NONE', '', 'NAN', 'NA'):
        score += 0.5
    else:
        # other weapons (e.g., blunt) medium
        score += 1.5

    # Crime domain (High weight)
    domain = str(row.get('Crime Domain', '')).upper()
    if domain in ['MURDER', 'HOMICIDE', 'RAPE', 'SEXUAL ASSAULT']:
        score += 3.0
    elif domain in ['ASSAULT', 'ROBBERY', 'BURGLARY']:
        score += 2.0
    elif domain in ['THEFT', 'HARASSMENT', 'CYBERCRIME']:
        score += 1.0
    else:
        score += 1.0

    # Victim age (Medium)
    age = row.get('Victim Age', None)
    try:
        age = float(age)
        if age < 18 or age > 60:
            score += 1.5
    except Exception:
        pass

    # Time of occurrence (Low)
    hour = row.get('Hour', None)
    try:
        hour = int(hour)
        # night = higher risk
        if (hour >= 20) or (hour <= 5):
            score += 1.0
    except Exception:
        pass

    # Location-based proxy: if the city has higher-than-average crime count add medium weight
    city = row.get('City', None)
    if city is not None and city in city_counts:
        if city_counts[city] > mean_city_count:
            score += 1.5

    # Cap & normalize to 1-10
    score = min(round(score, 2), 10.0)
    # Ensure minimum of 1
    if score < 1.0:
        score = 1.0
    return score

# Compute city crime counts for the filtered dataset (used by severity function)
if len(filtered_data) > 0 and 'City' in filtered_data.columns:
    city_counts_series = filtered_data['City'].value_counts()
    city_counts = city_counts_series.to_dict()
    mean_city_count = city_counts_series.mean()
else:
    city_counts = {}
    mean_city_count = 0

# Apply severity calculation
if len(filtered_data) > 0:
    filtered_data['Severity Index'] = filtered_data.apply(lambda r: calculate_severity(r, city_counts, mean_city_count), axis=1)
else:
    filtered_data['Severity Index'] = []

# -----------------------
# 4. Classification (XGBoost)
# -----------------------
if len(filtered_data) > 0 and 'Crime Domain_encoded' in filtered_data.columns and len(filtered_data['Crime Domain_encoded'].unique()) > 1:
    X_class = filtered_data[['City_encoded','Crime Code_encoded','Victim Age','Victim Gender_encoded',
                             'Weapon Used_encoded','Hour','Month']].fillna(0)
    y_class = filtered_data['Crime Domain_encoded']

    X_train, X_test, y_train, y_test = train_test_split(X_class, y_class, test_size=0.2, random_state=42)

    xgb_model = XGBClassifier(eval_metric='mlogloss', use_label_encoder=False)
    xgb_model.fit(X_train, y_train)
    y_pred = xgb_model.predict(X_test)

    st.subheader("XGBoost Classifier Accuracy")
    st.write("Accuracy:", accuracy_score(y_test, y_pred))

    cm = confusion_matrix(y_test, y_pred)
    try:
        classes = le_dict['Crime Domain'].classes_
        cm_df = pd.DataFrame(cm, index=classes, columns=classes)
    except Exception:
        cm_df = pd.DataFrame(cm)
    fig, ax = plt.subplots()
    sns.heatmap(cm_df, annot=True, fmt='d', ax=ax)
    st.pyplot(fig)
else:
    st.warning("Not enough data to train classifier with current filters.")

# -----------------------
# 5. Regression (HistGradientBoosting)
# -----------------------
if len(filtered_data) > 0 and 'Police Deployed' in filtered_data.columns:
    X_reg = filtered_data[['City_encoded','Crime Code_encoded','Victim Age','Victim Gender_encoded',
                          'Weapon Used_encoded','Hour','Month']].fillna(0)
    y_reg = filtered_data['Police Deployed'].fillna(0)

    scaler = StandardScaler()
    X_train_r, X_test_r, y_train_r, y_test_r = train_test_split(X_reg, y_reg, test_size=0.2, random_state=42)
    X_train_r = scaler.fit_transform(X_train_r)
    X_test_r = scaler.transform(X_test_r)

    hgb = HistGradientBoostingRegressor()
    hgb.fit(X_train_r, y_train_r)
    y_pred_r = hgb.predict(X_test_r)

    st.subheader("HistGradientBoosting Regression Results")
    st.write("Mean Squared Error:", mean_squared_error(y_test_r, y_pred_r))
else:
    st.warning("No data available for regression with current filters.")

# -----------------------
# 6. Clustering (KMeans)
# -----------------------
if len(filtered_data) > 0:
    cluster_features = [c for c in ['City_encoded','Crime Code_encoded','Victim Age','Hour'] if c in filtered_data.columns]
    if len(cluster_features) >= 2:
        X_cluster = filtered_data[cluster_features].fillna(0)
        scaler = StandardScaler()
        X_cluster_scaled = scaler.fit_transform(X_cluster)

        kmeans = KMeans(n_clusters=5, random_state=42, n_init=10)
        filtered_data['Cluster'] = kmeans.fit_predict(X_cluster_scaled)

        st.subheader("Crime Clusters (KMeans)")
        fig2 = px.scatter(
            filtered_data, x='Hour' if 'Hour' in filtered_data.columns else filtered_data.columns[0],
            y='Victim Age' if 'Victim Age' in filtered_data.columns else filtered_data.columns[0],
            color='Cluster',
            hover_data=['City','Crime Description'] if 'Crime Description' in filtered_data.columns else ['City']
        )
        st.plotly_chart(fig2)
    else:
        st.info("Not enough features for clustering.")
else:
    st.warning("No data available for clustering with current filters.")

# -----------------------
# 7. Visualizations (including Severity visuals)
# -----------------------
if len(filtered_data) > 0:
    # Severity visuals: show gauge for single-city selection, else show top-5 cities by avg severity
    st.subheader("⚠️ Crime Severity Index")

    # Average severity overall
    avg_sev_overall = filtered_data['Severity Index'].mean()

    # If user selected exactly one city — show that city's severity gauge
    if len(cities_selected) == 1 and cities_selected[0] in filtered_data['City'].unique():
        selected_city = cities_selected[0]
        city_avg = filtered_data[filtered_data['City'] == selected_city]['Severity Index'].mean()

        fig_gauge = go.Figure(go.Indicator(
            mode="gauge+number",
            value=round(city_avg, 2),
            title={'text': f"Severity Index — {selected_city}"},
            gauge={'axis': {'range': [0,10]},
                   'bar': {'color': "red"},
                   'steps': [
                       {'range': [0,3], 'color': "green"},
                       {'range': [3,6], 'color': "yellow"},
                       {'range': [6,10], 'color': "red"}]
                  }
        ))
        st.plotly_chart(fig_gauge)

    else:
        # Show overall gauge and a bar chart of top 5 cities by avg severity
        fig_gauge_all = go.Figure(go.Indicator(
            mode="gauge+number",
            value=round(avg_sev_overall, 2),
            title={'text': "Severity Index — Filtered Selection (Overall)"},
            gauge={'axis': {'range': [0,10]},
                   'bar': {'color': "red"},
                   'steps': [
                       {'range': [0,3], 'color': "green"},
                       {'range': [3,6], 'color': "yellow"},
                       {'range': [6,10], 'color': "red"}]
                  }
        ))
        st.plotly_chart(fig_gauge_all)

    # Continue with your other visuals...
    # Crimes per city
    st.subheader("Crime Trends by City")
    city_counts = filtered_data.groupby('City')['Crime Domain'].count().reset_index()
    fig3 = px.bar(city_counts, x='City', y='Crime Domain', title="Number of Crimes per City")
    st.plotly_chart(fig3)

    # Crimes per hour
    st.subheader("Crime Trends by Hour")
    hour_counts = filtered_data.groupby('Hour')['Crime Domain'].count().reset_index() if 'Hour' in filtered_data.columns else pd.DataFrame()
    if not hour_counts.empty:
        fig4 = px.line(hour_counts, x='Hour', y='Crime Domain', title="Crimes by Hour of Day")
        st.plotly_chart(fig4)

    # -----------------------
    # EXTRA VISUALIZATIONS
    # -----------------------

    # 1. Gender & Weapons
    if 'Victim Gender' in filtered_data.columns:
        st.subheader("Victim Gender Distribution")
        fig5 = px.pie(filtered_data, names="Victim Gender", title="Crimes by Victim Gender")
        st.plotly_chart(fig5)

    if 'Weapon Used' in filtered_data.columns:
        st.subheader("Weapons Used in Crimes")
        weapon_counts = filtered_data['Weapon Used'].value_counts().reset_index()
        weapon_counts.columns = ["Weapon", "Count"]
        fig6 = px.bar(weapon_counts, x="Weapon", y="Count", title="Weapons Used", color="Count")
        st.plotly_chart(fig6)

    # 2. Forecast (Prophet)
    if 'Date of Occurrence' in filtered_data.columns:
        st.subheader("🔮 Crime Forecast (Next 6 Months)")

        # Group data by month and count crimes
        monthly = filtered_data.groupby(filtered_data['Date of Occurrence'].dt.to_period('M'))['Crime Domain'].count().reset_index()
        monthly['Date'] = monthly['Date of Occurrence'].dt.to_timestamp()
        monthly = monthly.rename(columns={'Date': 'ds', 'Crime Domain': 'y'})

        # Train Prophet model (only if enough history)
        if len(monthly) >= 12:
            model = Prophet()
            model.fit(monthly)
            future = model.make_future_dataframe(periods=6, freq='M')
            forecast = model.predict(future)
            fig_forecast = px.line(forecast, x='ds', y='yhat', title="Forecasted Crimes for Next 6 Months")
            # highlight forecast region
            fig_forecast.add_vrect(
                x0=monthly['ds'].max(), x1=forecast['ds'].max(),
                fillcolor="lightblue", opacity=0.2,
                annotation_text="Forecasted Period", annotation_position="top left"
            )
            st.plotly_chart(fig_forecast)
        else:
            st.info("Not enough monthly history to run Prophet (need >=12 months).")

    # 4. Top 10 Crime Types
    st.subheader("Top 10 Crime Types")
    top_crimes = filtered_data['Crime Description'].value_counts().head(10).reset_index() if 'Crime Description' in filtered_data.columns else pd.DataFrame()
    if not top_crimes.empty:
        top_crimes.columns = ["Crime Description", "Count"]
        fig9 = px.bar(top_crimes, x="Crime Description", y="Count", title="Top 10 Crimes", color="Count")
        st.plotly_chart(fig9)

    # 5. Police Deployment vs Crime
    if 'Police Deployed' in filtered_data.columns:
        st.subheader("Police Deployment vs Crime Count")
        deploy_data = filtered_data.groupby('City')['Police Deployed'].sum().reset_index()
        fig10 = px.scatter(deploy_data, x="City", y="Police Deployed", size="Police Deployed", color="City",
                           title="Police Deployment by City")
        st.plotly_chart(fig10)

else:
    st.warning("No data available for visualizations with current filters.")

# -----------------------
# Alerts & Precautions (Bell Icon in Sidebar)
# -----------------------
st.sidebar.markdown("## 🔔 Notifications")

if st.sidebar.button("Show Alerts"):
    st.sidebar.subheader("🚨 Alerts & Precautions")

    if len(filtered_data) > 0:
        # 1. Highest crime city
        city_crime_counts = filtered_data['City'].value_counts()
        highest_city = city_crime_counts.idxmax()
        highest_count = city_crime_counts.max()
        st.sidebar.error(f"🔴 {highest_city} has the highest number of crimes ({highest_count})!")

        # 2. Most common crime
        if 'Crime Description' in filtered_data.columns:
            most_common_crime = filtered_data['Crime Description'].value_counts().idxmax()
            st.sidebar.warning(f"⚠️ Most common crime: {most_common_crime}")

            # Precautionary tips (example mapping, you can expand this)
            crime_precautions = {
                "THEFT": "Keep valuables secure and avoid isolated areas.",
                "ASSAULT": "Avoid late night travel alone and stay in well-lit places.",
                "BURGLARY": "Ensure home doors/windows are locked and use security cameras.",
                "ROBBERY": "Be cautious in crowded places and avoid displaying expensive items.",
            }
            if most_common_crime.upper() in crime_precautions:
                st.sidebar.info(f"💡 Precaution: {crime_precautions[most_common_crime.upper()]}")

        # 3. Trend alert (comparing current vs past months)
        if 'Date of Occurrence' in filtered_data.columns:
            monthly_trends = filtered_data.groupby(filtered_data['Date of Occurrence'].dt.to_period('M')).size()
            if len(monthly_trends) > 1:
                last, prev = monthly_trends.iloc[-1], monthly_trends.iloc[-2]
                if last > prev * 1.5:
                    st.sidebar.error("📈 Sharp increase in crimes compared to last month!")
                elif last < prev * 0.7:
                    st.sidebar.success("📉 Crimes have decreased significantly compared to last month!")
    else:
        st.sidebar.info("ℹ️ No data available for current filters.")
        

# -----------------------
# 9. Community & Social Insights (Community Awareness Only)
# -----------------------
st.subheader("🏘️ Community & Social Insights")

tab1 = st.tabs(["Community Awareness"])[0]

# -----------------------
# Tab 1: Community Awareness
# -----------------------
with tab1:
    st.write("🔹 **Community Awareness Recommendations**")

    if len(filtered_data) > 0 and 'Date of Occurrence' in filtered_data.columns:
        # Day, Hour, Month
        filtered_data['DayOfWeek'] = filtered_data['Date of Occurrence'].dt.day_name()
        day_counts = filtered_data.groupby('DayOfWeek')['Crime Domain'].count()
        hour_counts = filtered_data.groupby('Hour')['Crime Domain'].count() if 'Hour' in filtered_data.columns else pd.Series()
        month_counts = filtered_data.groupby('Month')['Crime Domain'].count() if 'Month' in filtered_data.columns else pd.Series()

        recommendations = []

        # Weekend tip
        if any(day_counts[day] > day_counts.mean() for day in ['Saturday','Sunday'] if day in day_counts.index):
            recommendations.append({
                "emoji": "🛡️",
                "title": "Weekend Theft Alert",
                "text": "Thefts increase on weekends. Secure your property and avoid carrying valuables unnecessarily."
            })

        # Evening/night tip
        if not hour_counts.empty:
            peak_hour = hour_counts.idxmax()
            if peak_hour >= 18:
                recommendations.append({
                    "emoji": "🌙",
                    "title": "Evening Crime Spike",
                    "text": f"Crime rates peak around {peak_hour}:00 hrs. Avoid isolated areas and travel in groups after dark."
                })

        # Seasonal tip (Summer May-Jun)
        if not month_counts.empty:
            summer_months = [5,6]
            summer_crimes = month_counts.loc[month_counts.index.isin(summer_months)].sum()
            if summer_crimes > month_counts.mean():
                recommendations.append({
                    "emoji": "☀️",
                    "title": "Summer Safety",
                    "text": "Higher crime rate during summer months. Stay alert, especially in crowded public areas."
                })

        # Display recommendations as info boxes
        if recommendations:
            for rec in recommendations:
                st.info(f"{rec['emoji']} **{rec['title']}**\n\n{rec['text']}")
        else:
            st.success("✅ No unusual patterns detected. Continue practicing regular safety measures.")
    else:
        st.warning("⚠️ Not enough data to generate community awareness tips.")
        
        
# -----------------------
# 10. Generate report options (PDF, Word, PPT)
# -----------------------

import streamlit as st
import zipfile
import io
from datetime import datetime
import matplotlib.pyplot as plt

# PDF
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER

# Word & PowerPoint
from docx import Document
from pptx import Presentation
from pptx.util import Inches

st.subheader("📄 Generate Auto Crime Reports")

# -----------------------
# PDF generation function (with charts)
# -----------------------
def generate_pdf(data, city_selection):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=18)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('TitleStyle', parent=styles['Title'], alignment=TA_CENTER, textColor=colors.darkblue)
    heading_style = ParagraphStyle('HeadingStyle', parent=styles['Heading2'], textColor=colors.darkred, spaceAfter=6)
    normal_style = ParagraphStyle('NormalStyle', parent=styles['Normal'], spaceAfter=4)

    elements = []

    # Title
    title = f"Crime Analysis Report — {', '.join(city_selection) if city_selection else 'All Cities'}"
    elements.append(Paragraph(title, title_style))
    elements.append(HRFlowable(width="100%", thickness=1, color=colors.black))
    elements.append(Spacer(1, 12))

    # Summary
    elements.append(Paragraph("📊 Summary of Analysis", heading_style))
    total_crimes = len(data)
    avg_severity = round(data['Severity Index'].mean(), 2) if 'Severity Index' in data else "N/A"
    elements.append(Paragraph(f"• Total Reported Crimes: <b>{total_crimes}</b>", normal_style))
    elements.append(Paragraph(f"• Average Severity Index: <b>{avg_severity}</b>", normal_style))
    elements.append(Spacer(1, 12))

    # Top Cities Table & Chart
    if 'City' in data.columns:
        city_stats = data['City'].value_counts().head(5)
        city_data_table = [["City", "Crime Count"]] + [[c, int(n)] for c, n in city_stats.items()]
        table = Table(city_data_table, hAlign='LEFT', colWidths=[150, 100])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.darkblue),
            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.lightgrey, colors.whitesmoke]),
            ('GRID',(0,0),(-1,-1),0.5,colors.black)
        ]))
        elements.append(Paragraph("🏙️ Top Cities by Crime Count", heading_style))
        elements.append(table)
        elements.append(Spacer(1, 12))

        # Bar chart
        fig, ax = plt.subplots()
        ax.bar(city_stats.index, city_stats.values, color='skyblue')
        ax.set_title("Top Cities by Crime Count")
        ax.set_ylabel("Crime Count")
        ax.set_xlabel("City")
        plt.tight_layout()
        chart_buffer = io.BytesIO()
        plt.savefig(chart_buffer, format='PNG')
        plt.close(fig)
        chart_buffer.seek(0)
        elements.append(Image(chart_buffer, width=400, height=250))
        elements.append(Spacer(1, 12))

    # Common Crime Pattern
    if 'Crime Description' in data.columns:
        top_crime = data['Crime Description'].value_counts().idxmax()
        count = data['Crime Description'].value_counts().max()
        elements.append(Paragraph("🔍 Common Crime Pattern", heading_style))
        elements.append(Paragraph(f"Most common crime: <b>{top_crime}</b> ({count} cases)", normal_style))
        elements.append(Spacer(1, 12))

    # Alerts
    elements.append(Paragraph("🚨 Alerts & Trends", heading_style))
    if 'Date of Occurrence' in data.columns:
        monthly_trends = data.groupby(data['Date of Occurrence'].dt.to_period('M')).size()
        if len(monthly_trends) > 1:
            last, prev = monthly_trends.iloc[-1], monthly_trends.iloc[-2]
            if last > prev * 1.5:
                elements.append(Paragraph("⚠️ Sharp increase in crimes compared to last month!", normal_style))
            elif last < prev * 0.7:
                elements.append(Paragraph("✅ Crimes have decreased significantly compared to last month.", normal_style))
            else:
                elements.append(Paragraph("ℹ️ Crime rates are stable compared to last month.", normal_style))
    elements.append(Spacer(1, 12))

    # Recommendations
    elements.append(Paragraph("📌 Strategic Recommendations", heading_style))
    recommendations = [
        "Increase patrols in high-severity areas.",
        "Install CCTV and improve street lighting.",
        "Conduct community awareness programs.",
        "Deploy gender-sensitive safety units.",
        "Use predictive analytics for proactive policing."
    ]
    for rec in recommendations:
        elements.append(Paragraph(f"• {rec}", normal_style))
    elements.append(Spacer(1, 12))

    # Footer
    now = datetime.now().strftime("%d-%m-%Y %H:%M")
    elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.grey))
    elements.append(Paragraph(f"<i>Generated automatically by Crime Analysis Dashboard on {now}</i>", styles['Normal']))

    doc.build(elements)
    buffer.seek(0)
    return buffer

# -----------------------
# Word generation function
# -----------------------
def generate_word(data, city_selection):
    doc = Document()
    doc.add_heading(f"Crime Analysis Report — {', '.join(city_selection) if city_selection else 'All Cities'}", 0)

    doc.add_heading("Summary of Analysis", level=1)
    total_crimes = len(data)
    avg_severity = round(data['Severity Index'].mean(), 2) if 'Severity Index' in data else "N/A"
    doc.add_paragraph(f"• Total Reported Crimes: {total_crimes}")
    doc.add_paragraph(f"• Average Severity Index: {avg_severity}")

    # Top Cities
    if 'City' in data.columns:
        city_stats = data['City'].value_counts().head(5)
        doc.add_heading("Top Cities by Crime Count", level=2)
        table = doc.add_table(rows=1, cols=2)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'City'
        hdr_cells[1].text = 'Crime Count'
        for city, count in city_stats.items():
            row_cells = table.add_row().cells
            row_cells[0].text = city
            row_cells[1].text = str(count)

    # Save to BytesIO
    word_buffer = io.BytesIO()
    doc.save(word_buffer)
    word_buffer.seek(0)
    return word_buffer

# -----------------------
# PowerPoint generation function
# -----------------------
def generate_ppt(data, city_selection):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank slide

    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = f"Crime Analysis Report — {', '.join(city_selection) if city_selection else 'All Cities'}"

    # Summary
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text_frame.text = "Summary of Analysis"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3)).text_frame
    total_crimes = len(data)
    avg_severity = round(data['Severity Index'].mean(), 2) if 'Severity Index' in data else "N/A"
    tf.text = f"• Total Reported Crimes: {total_crimes}\n• Average Severity Index: {avg_severity}"

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# -----------------------
# Streamlit buttons
# -----------------------
# Single PDF
if st.button("📄 Generate PDF Report (All Cities)"):
    if len(filtered_data) > 0:
        pdf_buffer = generate_pdf(filtered_data, cities_selected)
        st.download_button(
            label="⬇️ Download PDF",
            data=pdf_buffer,
            file_name="Crime_Analysis_Report.pdf",
            mime="application/pdf"
        )

# PDF ZIP per city
if st.button("📄 Generate PDF Reports per City (ZIP)"):
    if len(filtered_data) > 0 and len(cities_selected) > 0:
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w") as zip_file:
            for city in cities_selected:
                city_data = filtered_data[filtered_data['City'] == city]
                pdf_buf = generate_pdf(city_data, [city])
                zip_file.writestr(f"Crime_Report_{city}.pdf", pdf_buf.getvalue())
        zip_buffer.seek(0)
        st.download_button(
            label="⬇️ Download All City Reports (ZIP)",
            data=zip_buffer,
            file_name="Crime_Reports_Per_City.zip",
            mime="application/zip"
        )

# Word
if st.button("📝 Generate Word Report (All Cities)"):
    if len(filtered_data) > 0:
        word_buffer = generate_word(filtered_data, cities_selected)
        st.download_button(
            label="⬇️ Download Word",
            data=word_buffer,
            file_name="Crime_Analysis_Report.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

# PowerPoint
if st.button("📊 Generate PPT Report (All Cities)"):
    if len(filtered_data) > 0:
        ppt_buffer = generate_ppt(filtered_data, cities_selected)
        st.download_button(
            label="⬇️ Download PPT",
            data=ppt_buffer,
            file_name="Crime_Analysis_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if st.session_state.chat_started:

    st.subheader("💬 Crime AI Assistant")

    # Show past chat
    for msg in st.session_state.messages:
        if msg["role"] == "assistant":
            st.write("🤖:", msg["content"])
        else:
            st.write("🧑‍💻:", msg["content"])

    # User input
    user_query = st.text_input("Your question:")

    if st.button("Send"):
        if user_query.strip():
            st.session_state.messages.append({"role": "user", "content": user_query})

            ai_answer = get_chat_response(user_query)
            st.session_state.messages.append({"role": "assistant", "content": ai_answer})

            st.experimental_rerun()
